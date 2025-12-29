import streamlit as st
import os
import zipfile
import json
import re
import uuid
import socket
import io
import posixpath
import xml.etree.ElementTree as ET
from typing import Optional, List, Tuple, Set

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from googleapiclient.errors import HttpError
from google.auth.transport.requests import Request

# --- 設定全域超時 (100分鐘) ---
socket.setdefaulttimeout(6000)

SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/presentations",
    "https://www.googleapis.com/auth/spreadsheets",
]

# 請確認這是你的正確 Google Sheet ID
SPREADSHEET_ID = "1tkLPKqFQld2bCythqNY0CX83w4y1cWZJvW6qErE8vek"
# 固定權限管理員字串
PERMITTED_ADMINS_STRING = "admin,william,robot,fm,sunny,jason,eq,com,mona"

VIDEO_EXTS = (".mp4", ".mov", ".avi", ".m4v", ".wmv")
IMAGE_EXTS = (".jpg", ".jpeg", ".png", ".tiff", ".bmp")

# Namespaces / rel types
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
OFFICE_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

OFFICE_DOC_REL = f"{OFFICE_NS}/officeDocument"
SLIDE_REL_TYPE = f"{OFFICE_NS}/slide"


def natural_sort_key(s: str):
    return [int(text) if text.isdigit() else text.lower()
            for text in re.split(r"([0-9]+)", s)]


def _log(log_callback, msg: str):
    if log_callback:
        log_callback(msg)


def _normalize_part_path(path: str) -> str:
    path = path.replace("\\", "/").lstrip("/")
    return posixpath.normpath(path)


def _rels_path_for_part(part_path: str) -> str:
    part_path = _normalize_part_path(part_path)
    base_dir = posixpath.dirname(part_path)
    filename = posixpath.basename(part_path)
    return posixpath.join(base_dir, "_rels", f"{filename}.rels")


def _resolve_target(base_part: str, target: str) -> str:
    target = (target or "").replace("\\", "/")
    if target.startswith("/"):
        return _normalize_part_path(target)
    base_dir = posixpath.dirname(_normalize_part_path(base_part))
    return _normalize_part_path(posixpath.join(base_dir, target))


def _read_from_zip(z: zipfile.ZipFile, name: str) -> Optional[bytes]:
    try:
        return z.read(name)
    except KeyError:
        return None


def _is_external_rel(rel_el: ET.Element) -> bool:
    return rel_el.attrib.get("TargetMode", "").lower() == "external"


def _parse_relationship_targets(rels_xml: bytes) -> List[Tuple[str, bool]]:
    ns = {"r": PKG_REL_NS}
    root = ET.fromstring(rels_xml)
    out: List[Tuple[str, bool]] = []
    for rel in root.findall("r:Relationship", ns):
        out.append((rel.attrib.get("Target", ""), _is_external_rel(rel)))
    return out


def _strip_video_relationships(rels_xml: bytes) -> bytes:
    try:
        ns = {"r": PKG_REL_NS}
        root = ET.fromstring(rels_xml)
        changed = False

        for rel in list(root.findall("r:Relationship", ns)):
            if _is_external_rel(rel):
                continue
            target = (rel.attrib.get("Target", "") or "").lower().replace("\\", "/")
            if target.endswith(VIDEO_EXTS) and (
                "/media/" in target or target.startswith("../media/") or target.startswith("media/")
            ):
                root.remove(rel)
                changed = True

        if not changed:
            return rels_xml
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    except Exception:
        return rels_xml


def _ensure_officedocument_in_root_rels(root_rels_xml: Optional[bytes]) -> bytes:
    if not root_rels_xml:
        root = ET.Element("Relationships", xmlns=PKG_REL_NS)
        ET.SubElement(root, "Relationship", {
            "Id": "rId1",
            "Type": OFFICE_DOC_REL,
            "Target": "ppt/presentation.xml",
        })
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    try:
        ns = {"r": PKG_REL_NS}
        root = ET.fromstring(root_rels_xml)

        for rel in root.findall("r:Relationship", ns):
            if rel.attrib.get("Type") == OFFICE_DOC_REL:
                return root_rels_xml

        ET.SubElement(root, "Relationship", {
            "Id": "rId1",
            "Type": OFFICE_DOC_REL,
            "Target": "ppt/presentation.xml",
        })
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    except Exception:
        root = ET.Element("Relationships", xmlns=PKG_REL_NS)
        ET.SubElement(root, "Relationship", {
            "Id": "rId1",
            "Type": OFFICE_DOC_REL,
            "Target": "ppt/presentation.xml",
        })
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _used_slide_rids_from_presentation_xml(presentation_xml: bytes) -> Set[str]:
    ns = {"p": PML_NS, "r": OFFICE_NS}
    root = ET.fromstring(presentation_xml)
    used: Set[str] = set()

    for sldId in root.findall(".//p:sldIdLst/p:sldId", ns):
        rid = sldId.attrib.get(f"{{{OFFICE_NS}}}id")
        if rid:
            used.add(rid)

    return used


def _rebuild_presentation_rels(pres_rels_xml: bytes, used_slide_rids: Set[str]) -> bytes:
    ns = {"r": PKG_REL_NS}
    root = ET.fromstring(pres_rels_xml)
    changed = False

    for rel in list(root.findall("r:Relationship", ns)):
        rel_type = rel.attrib.get("Type", "")
        rel_id = rel.attrib.get("Id", "")
        if rel_type == SLIDE_REL_TYPE and rel_id not in used_slide_rids:
            root.remove(rel)
            changed = True

    if not changed:
        return pres_rels_xml
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _prune_content_types_overrides(ct_xml: bytes, keep_parts: Set[str]) -> bytes:
    try:
        ns = {"ct": CT_NS}
        root = ET.fromstring(ct_xml)
        keep_with_slash = {"/" + _normalize_part_path(p) for p in keep_parts}

        changed = False
        for override in list(root.findall("ct:Override", ns)):
            part_name = override.attrib.get("PartName", "")
            if part_name and part_name not in keep_with_slash:
                root.remove(override)
                changed = True

        if not changed:
            return ct_xml
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)
    except Exception:
        return ct_xml


class PPTAutomationBot:
    def __init__(self):
        self.creds = self._get_credentials()
        
        if self.creds:
            self.drive_service = build("drive", "v3", credentials=self.creds)
            self.slides_service = build("slides", "v1", credentials=self.creds)
            self.sheets_service = build("sheets", "v4", credentials=self.creds)
        else:
            self.drive_service = None
            self.slides_service = None
            self.sheets_service = None

    def _get_credentials(self):
        creds = None
        if "google_token" in st.secrets:
            try:
                token_info = json.loads(st.secrets["google_token"])
                creds = Credentials.from_authorized_user_info(token_info, SCOPES)
            except Exception as e:
                print(f"雲端 Token 讀取失敗: {e}")

        if not creds and os.path.exists('token.json'):
            try:
                creds = Credentials.from_authorized_user_file('token.json', SCOPES)
            except Exception as e:
                print(f"本機 token.json 讀取失敗: {e}")

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                try:
                    creds.refresh(Request())
                except Exception as e:
                    st.error(f"Token 過期且無法自動刷新: {e}")
                    return None
            else:
                st.error("找不到有效的憑證，請確認已設定 Streamlit Secrets。")
                return None
        
        return creds

    def get_user_email(self):
        if not self.drive_service:
            return "服務未初始化"
        try:
            about = self.drive_service.about().get(fields="user").execute()
            return about["user"]["emailAddress"]
        except Exception:
            return "未知"

    def _check_drive_file_exists(self, filename):
        try:
            query = f"name = '{filename}' and trashed = false"
            results = self.drive_service.files().list(
                q=query, spaces="drive", fields="files(id, name, webViewLink)"
            ).execute()
            files = results.get("files", [])
            if files:
                return files[0].get("id"), files[0].get("webViewLink")
        except Exception as e:
            print(f"查詢 Drive 失敗: {e}")
        return None

    def _create_play_icon(self, filename):
        if os.path.exists(filename):
            return
        # 創建一個簡單的播放圖示 (灰色背景)
        img = Image.new("RGB", (200, 150), color=(100, 100, 100))
        img.save(filename)

    # =========================
    #  核心：拆分後清理邏輯
    # =========================
    def _prune_pptx_package_fast(self, pptx_path: str, log_callback=None) -> None:
        if not os.path.exists(pptx_path):
            return

        tmp_out = f"{pptx_path}.pruned_{uuid.uuid4().hex[:6]}.pptx"

        with zipfile.ZipFile(pptx_path, "r") as zin:
            names = set(zin.namelist())

            root_rels_name = "_rels/.rels"
            root_rels_xml = _read_from_zip(zin, root_rels_name)
            fixed_root_rels = _ensure_officedocument_in_root_rels(root_rels_xml)

            pres_xml_name = "ppt/presentation.xml"
            pres_rels_name = "ppt/_rels/presentation.xml.rels"

            pres_xml = _read_from_zip(zin, pres_xml_name)
            pres_rels_xml = _read_from_zip(zin, pres_rels_name)

            pres_rels_fixed: Optional[bytes] = None
            if pres_xml and pres_rels_xml:
                used_slide_rids = _used_slide_rids_from_presentation_xml(pres_xml)
                pres_rels_clean = _strip_video_relationships(pres_rels_xml)
                pres_rels_fixed = _rebuild_presentation_rels(pres_rels_clean, used_slide_rids)

            def get_rels_xml(rels_name: str) -> Optional[bytes]:
                b = _read_from_zip(zin, rels_name)
                if b is None:
                    return None
                if rels_name == pres_rels_name and pres_rels_fixed is not None:
                    return pres_rels_fixed
                return _strip_video_relationships(b)

            keep: Set[str] = set()
            keep.add("[Content_Types].xml")
            keep.add(root_rels_name)

            queue: List[str] = []

            try:
                for target, is_ext in _parse_relationship_targets(fixed_root_rels):
                    if is_ext:
                        continue
                    resolved = _resolve_target("", target)
                    if resolved in names:
                        if resolved.startswith("ppt/media/") and resolved.lower().endswith(VIDEO_EXTS):
                            continue
                        queue.append(resolved)
            except Exception:
                pass

            while queue:
                part = _normalize_part_path(queue.pop())
                if part in keep:
                    continue
                if part.startswith("ppt/media/") and part.lower().endswith(VIDEO_EXTS):
                    continue
                if part not in names:
                    continue

                keep.add(part)

                rels_name = _rels_path_for_part(part)
                if rels_name in names:
                    keep.add(rels_name)
                    rels_xml = get_rels_xml(rels_name)
                    if rels_xml:
                        try:
                            for target, is_ext in _parse_relationship_targets(rels_xml):
                                if is_ext:
                                    continue
                                resolved = _resolve_target(part, target)
                                if resolved in names:
                                    if resolved.startswith("ppt/media/") and resolved.lower().endswith(VIDEO_EXTS):
                                        continue
                                    queue.append(resolved)
                        except Exception:
                            pass

            for maybe in ("docProps/app.xml", "docProps/core.xml"):
                if maybe in names:
                    keep.add(maybe)
                    rels_name = _rels_path_for_part(maybe)
                    if rels_name in names:
                        keep.add(rels_name)

            with zipfile.ZipFile(tmp_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                ct_xml = _read_from_zip(zin, "[Content_Types].xml")
                if ct_xml:
                    ct_xml2 = _prune_content_types_overrides(ct_xml, keep)
                    zout.writestr("[Content_Types].xml", ct_xml2)

                zout.writestr(root_rels_name, fixed_root_rels)

                for name in sorted(keep):
                    if name in ("[Content_Types].xml", root_rels_name):
                        continue
                    if name.startswith("ppt/media/") and name.lower().endswith(VIDEO_EXTS):
                        continue
                    if name not in names:
                        continue

                    if name.lower().endswith(".rels"):
                        b = get_rels_xml(name)
                        if b is None:
                            continue
                        zout.writestr(name, b)
                    elif name == pres_rels_name and pres_rels_fixed is not None:
                        zout.writestr(name, pres_rels_fixed)
                    else:
                        zout.writestr(name, zin.read(name))

        os.replace(tmp_out, pptx_path)
        _log(log_callback, f"✅ [Prune] {os.path.basename(pptx_path)}：清理完成。")

    # === Step 1: 提取與上傳影片 ===
    def extract_and_upload_videos(self, pptx_path, extract_dir, file_prefix="", progress_callback=None, log_callback=None):
        if not self.drive_service:
            _log(log_callback, "❌ 服務未初始化，無法上傳影片。")
            return {}

        if not os.path.exists(extract_dir):
            os.makedirs(extract_dir)

        safe_prefix = file_prefix if file_prefix else "default"
        map_filename = f"video_map_{safe_prefix}.json"
        map_path = map_filename

        video_map = {}
        if os.path.exists(map_path):
            try:
                with open(map_path, "r", encoding="utf-8") as f:
                    video_map = json.load(f)
            except Exception:
                pass

        with zipfile.ZipFile(pptx_path, "r") as z:
            video_files = [
                f for f in z.infolist()
                if f.filename.startswith("ppt/media/")
                and f.filename.lower().endswith(VIDEO_EXTS)
            ]
            video_files.sort(key=lambda f: natural_sort_key(os.path.basename(f.filename)))
            total_videos = len(video_files)

            _log(log_callback, f"📊 掃描完成：共發現 {total_videos} 個影片檔。")

            for idx, file_info in enumerate(video_files):
                original_filename = os.path.basename(file_info.filename)

                if original_filename in video_map:
                    _log(log_callback, f"⏭️ ({idx+1}/{total_videos}) {original_filename} 本地紀錄已存在，跳過。")
                    continue

                _log(log_callback, f"📦 ({idx+1}/{total_videos}) 正在解壓縮與查重：{original_filename} ...")

                z.extract(file_info, extract_dir)
                full_path = os.path.join(extract_dir, file_info.filename)

                upload_name = f"[{file_prefix}]_{original_filename}" if file_prefix else original_filename

                existing_file = self._check_drive_file_exists(upload_name)
                if existing_file:
                    _, web_link = existing_file
                    _log(log_callback, f"☁️ ({idx+1}/{total_videos}) 雲端已有檔案：{upload_name}，直接使用！")
                    video_map[original_filename] = web_link
                    with open(map_path, "w", encoding="utf-8") as f:
                        json.dump(video_map, f, indent=4)
                    continue

                _log(log_callback, f"⬆️ ({idx+1}/{total_videos}) 開始上傳：{upload_name} ...")

                try:
                    file_metadata = {"name": upload_name}
                    CHUNK_SIZE = 5 * 1024 * 1024
                    media = MediaFileUpload(full_path, resumable=True, chunksize=CHUNK_SIZE)

                    request = self.drive_service.files().create(
                        body=file_metadata, media_body=media, fields="id, webViewLink"
                    )

                    response = None
                    while response is None:
                        status, response = request.next_chunk()
                        if status and progress_callback:
                            # 這裡傳遞的是單個檔案的上傳進度
                            progress_callback(upload_name, int(status.resumable_progress), int(status.total_size))

                    file = response
                    self.drive_service.permissions().create(
                        fileId=file.get("id"),
                        body={"type": "anyone", "role": "reader"},
                    ).execute()

                    video_map[original_filename] = file.get("webViewLink")
                    with open(map_path, "w", encoding="utf-8") as f:
                        json.dump(video_map, f, indent=4)

                except Exception as e:
                    print(f"上傳失敗: {e}")
                    pass

        return video_map

    # === Step 2: 置換為圖片連結 (加入進度回報) ===
    def replace_videos_with_images(self, input_pptx, output_pptx, video_map, progress_callback=None):
        if os.path.exists(output_pptx):
            print(f"Step 2: {output_pptx} 已存在，跳過。")
            return

        icon_path = "play_icon.png"
        self._create_play_icon(icon_path)

        prs = Presentation(input_pptx)
        total_slides = len(prs.slides) # 計算總頁數用於進度

        for i, slide in enumerate(prs.slides):
            # 回報進度
            if progress_callback:
                progress_callback(i + 1, total_slides)

            slide_video_filenames = []
            for rel in slide.part.rels.values():
                if "media" in rel.target_ref:
                    fname = os.path.basename(rel.target_ref)
                    if fname in video_map:
                        slide_video_filenames.append(fname)

            shapes_to_replace = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    target_filename = None
                    if len(slide_video_filenames) >= 1:
                        target_filename = slide_video_filenames[0]

                    if target_filename and target_filename in video_map:
                        shapes_to_replace.append({
                            "shape": shape,
                            "link": video_map[target_filename],
                            "left": shape.left, "top": shape.top,
                            "width": shape.width, "height": shape.height,
                        })

            for item in shapes_to_replace:
                sp = item["shape"]
                sp.element.getparent().remove(sp.element)
                pic = slide.shapes.add_picture(
                    icon_path, item["left"], item["top"], item["width"], item["height"]
                )
                pic.click_action.hyperlink.address = item["link"]

        prs.save(output_pptx)

    # === Step 3: 檔案瘦身 (加入進度回報) ===
    def shrink_pptx(self, input_pptx, output_pptx, progress_callback=None):
        if os.path.exists(output_pptx):
            print(f"Step 3: {output_pptx} 已存在，跳過。")
            return

        print("🚀 開始執行 Step 3: 圖片壓縮 (1280px/Q50)...")

        with zipfile.ZipFile(input_pptx, "r") as zin:
            # 計算總檔案數用於進度
            file_list = zin.infolist()
            total_files = len(file_list)

            with zipfile.ZipFile(output_pptx, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                for i, item in enumerate(file_list):
                    # 回報進度
                    if progress_callback:
                        progress_callback(i + 1, total_files)

                    name = item.filename

                    # 移除影片實體
                    if name.startswith("ppt/media/") and name.lower().endswith(VIDEO_EXTS):
                        continue

                    # 處理圖片
                    if name.startswith("ppt/media/") and name.lower().endswith(IMAGE_EXTS):
                        try:
                            file_data = zin.read(name)
                            # 小於 50KB 不壓縮
                            if len(file_data) < 50 * 1024:
                                zout.writestr(item, file_data)
                                continue

                            img = Image.open(io.BytesIO(file_data))
                            
                            # [規格] 1280px
                            img.thumbnail((1280, 1280), Image.Resampling.LANCZOS)

                            output_buffer = io.BytesIO()
                            ext = os.path.splitext(name)[1].lower()

                            # [規格] Quality 50
                            if ext in (".jpg", ".jpeg"):
                                img = img.convert("RGB")
                                img.save(output_buffer, format="JPEG", quality=50, optimize=True)
                                zout.writestr(name, output_buffer.getvalue())
                            elif ext == ".png":
                                img.save(output_buffer, format="PNG", optimize=True)
                                zout.writestr(name, output_buffer.getvalue())
                            else:
                                zout.writestr(item, file_data)
                            continue

                        except Exception as e:
                            print(f"   ❌ 壓縮圖片 {name} 失敗: {e}，保留原圖。")
                            zout.writestr(item, zin.read(name))
                            continue

                    zout.writestr(item, zin.read(name))

    # === Step 4: 拆分與上傳 (加入前綴處理) ===
    def split_and_upload(self, slim_pptx, split_jobs, file_prefix="", progress_callback=None, log_callback=None, debug_mode=False):
        if not self.drive_service:
            _log(log_callback, "❌ 服務未初始化，無法上傳拆分檔。")
            return []

        results = []
        total_jobs = len(split_jobs)

        # Debug 模式目錄 (如果未來需要啟用)
        debug_dir = "debug_output"
        if debug_mode and not os.path.exists(debug_dir):
            os.makedirs(debug_dir)

        for idx, job in enumerate(split_jobs):
            current_num = idx + 1
            original_filename = job["filename"]
            # [新增] 加上前綴的最終顯示檔名
            display_name = f"[{file_prefix}]_{original_filename}" if file_prefix else original_filename
            
            # 確保副檔名
            if not display_name.endswith('.pptx'):
                 display_name += ".pptx"

            # Debug Mode (略)
            if debug_mode:
                results.append(job)
                continue

            # 一般模式
            if job.get("final_link"):
                _log(log_callback, f"⏭️ ({current_num}/{total_jobs}) {display_name} 本地已完成，跳過。")
                results.append(job)
                continue

            existing_file = self._check_drive_file_exists(display_name)
            if existing_file:
                file_id, web_link = existing_file
                _log(log_callback, f"☁️ ({current_num}/{total_jobs}) 雲端已有簡報：{display_name}，直接使用！")
                job["final_link"] = web_link
                job["presentation_id"] = file_id
                results.append(job)
                continue

            temp_split_name = f"temp_{uuid.uuid4().hex[:6]}.pptx"
            try:
                _log(log_callback, f"✂️ ({current_num}/{total_jobs}) 正在拆分：{display_name} ...")

                prs = Presentation(slim_pptx)
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                # 轉換為 0-based index
                keep_indices = set(range(job["start"] - 1, job["end"]))

                # 倒序刪除不需要的投影片
                for i in range(len(slides) - 1, -1, -1):
                    if i not in keep_indices:
                        xml_slides.remove(slides[i])

                prs.save(temp_split_name)

                # 拆分後立刻執行清理
                try:
                    self._prune_pptx_package_fast(temp_split_name, log_callback=log_callback)
                except Exception as e:
                    _log(log_callback, f"⚠️ [Prune] 失敗但不致命：{e}")

                file_size = os.path.getsize(temp_split_name)
                size_mb = file_size / (1024 * 1024)

                if size_mb > 99:
                    error_msg = f"⛔️ 檔案過大：{display_name} 仍有 {size_mb:.2f} MB (超過 100MB 限制)。"
                    _log(log_callback, error_msg)
                    job["error_too_large"] = True
                    job["size_mb"] = size_mb
                    results.append(job)
                    continue

                _log(log_callback, f"⬆️ ({current_num}/{total_jobs}) 正在上傳：{display_name} (大小: {size_mb:.2f} MB)...")

                file_metadata = {"name": display_name, "mimeType": "application/vnd.google-apps.presentation"}

                CHUNK_SIZE = 5 * 1024 * 1024
                media = MediaFileUpload(
                    temp_split_name,
                    mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    resumable=True,
                    chunksize=CHUNK_SIZE,
                )

                request = self.drive_service.files().create(
                    body=file_metadata, media_body=media, fields="id, webViewLink"
                )

                response = None
                while response is None:
                    status, response = request.next_chunk()
                    if status and progress_callback:
                        # 回報單檔上傳進度
                        progress_callback(display_name, int(status.resumable_progress), int(status.total_size))

                file = response
                self.drive_service.permissions().create(
                    fileId=file.get("id"), body={"type": "anyone", "role": "reader"}
                ).execute()

                job["final_link"] = file.get("webViewLink")
                job["presentation_id"] = file.get("id")
                results.append(job)

            except Exception as e:
                print(f"上傳失敗: {e}")
                results.append(job)
            finally:
                if os.path.exists(temp_split_name):
                    os.remove(temp_split_name)

        return results

    # === Step 5: 內嵌優化 (加入進度回報) ===
    def embed_videos_in_slides(self, processed_jobs, progress_callback=None, log_callback=None, debug_mode=False):
        if debug_mode:
            return processed_jobs
        
        if not self.slides_service:
            return processed_jobs

        jobs_to_process = [j for j in processed_jobs if "presentation_id" in j]
        total_jobs = len(jobs_to_process)
        count = 0

        for job in jobs_to_process:
            count += 1
            # 回報進度
            if progress_callback:
                progress_callback(count, total_jobs)

            pid = job["presentation_id"]
            _log(log_callback, f"🔧 ({count}/{total_jobs}) 正在優化播放器：{job['filename']} ...")

            try:
                presentation = self.slides_service.presentations().get(presentationId=pid).execute()
                requests = []

                for slide in presentation.get("slides", []):
                    page_id = slide["objectId"]
                    for element in slide.get("pageElements", []):
                        if "image" in element:
                            url = element["image"].get("imageProperties", {}).get("link", {}).get("url", "")
                            if "drive.google.com" in url:
                                match = re.search(r"/file/d/([a-zA-Z0-9-_]+)", url)
                                if match:
                                    vid_id = match.group(1)
                                    requests.append({
                                        "createVideo": {
                                            "source": "DRIVE",
                                            "id": vid_id,
                                            "elementProperties": {
                                                "pageObjectId": page_id,
                                                "size": element.get("size"),
                                                "transform": element.get("transform"),
                                            },
                                        }
                                    })
                                    requests.append({"deleteObject": {"objectId": element["objectId"]}})

                if requests:
                    self.slides_service.presentations().batchUpdate(
                        presentationId=pid, body={"requests": requests}
                    ).execute()

            except Exception as e:
                print(f"優化失敗: {e}")

        return processed_jobs

    # === Step 6: 寫入 Google Sheet (欄位調整) ===
    def log_to_sheets(self, completed_jobs, log_callback=None, debug_mode=False):
        if debug_mode:
            return
        
        if not self.sheets_service:
            _log(log_callback, "❌ 服務未初始化，無法寫入試算表。")
            return

        existing_ids = set()
        try:
            _log(log_callback, "🔍 正在比對 Google Sheet 既有資料，避免重複寫入...")
            # 讀取 A 欄比對 ID
            result = self.sheets_service.spreadsheets().values().get(
                spreadsheetId=SPREADSHEET_ID,
                range="Presentations!A:A",
            ).execute()
            rows = result.get("values", [])
            for r in rows:
                if r:
                    existing_ids.add(r[0])

        except HttpError as err:
            if err.resp.status == 403:
                current_email = self.get_user_email()
                error_msg = (
                    f"⛔️ 權限錯誤 (403)：請將表單分享給此 Email 並設為「編輯者」：\n"
                    f"👉 {current_email}"
                )
                print(error_msg)
                raise Exception(error_msg)
            raise err

        except Exception as e:
            print(f"讀取 Sheet 失敗: {e}")
            raise e

        values = []
        jobs_to_mark_done = []

        for job in completed_jobs:
            if "final_link" not in job:
                continue

            job_id = job.get("id")
            if job_id in existing_ids:
                _log(log_callback, f"⏭️ 任務 {job['filename']} (ID: {job_id}) 已存在於報表中，跳過。")
                continue

            # [修正] 欄位順序: 
            # id, Category, SubCategory, Region, Client, SlideURL, Keywords, title, PermittedAdmins, CustomThumbnail
            row = [
                job_id,
                job.get("category", ""),
                job.get("subcategory", ""),
                "", # Region (目前無此欄位，填空)
                job.get("client", ""),
                job["final_link"], # SlideURL
                job.get("keywords", ""),
                job["filename"], # title (檔名)
                PERMITTED_ADMINS_STRING, # PermittedAdmins
                ""  # CustomThumbnail (目前無縮圖，填空)
            ]
            values.append(row)
            jobs_to_mark_done.append(job)

        if values:
            _log(log_callback, f"📝 正在寫入 {len(values)} 筆新資料到 Google Sheets...")

            body = {"values": values}
            self.sheets_service.spreadsheets().values().append(
                spreadsheetId=SPREADSHEET_ID,
                range="Presentations!A:J", # [修正] 範圍擴大到 J 欄
                valueInputOption="USER_ENTERED",
                body=body,
            ).execute()

            for job in jobs_to_mark_done:
                job["logged_to_sheet"] = True
        else:
            _log(log_callback, "✅ 所有資料皆已存在於 Sheet 中，同步完成。")