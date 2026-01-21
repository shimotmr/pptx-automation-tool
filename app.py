# Version: v1.4 (Diagnostic Mode)
# Focus: Catching Step 2 (Video Replace) errors and forcing display on Sidebar.

import streamlit as st
import streamlit.components.v1 as components
import os
import uuid
import json
import shutil
import traceback
import requests
from pptx import Presentation

# --- [診斷 1] 強制檢查必要套件 ---
try:
    import PIL
    import lxml
except ImportError as e:
    st.error(f"❌ 嚴重環境錯誤：缺少必要套件！\n請在 requirements.txt 加入：\nPillow\nlxml\n\n詳細錯誤: {e}")
    st.stop()

# 嘗試載入處理器
try:
    from ppt_processor import PPTAutomationBot
except ImportError:
    st.sidebar.error("❌ 找不到 `ppt_processor.py`，請確認檔案已上傳！")
    st.stop()

# ==========================================
#              設定頁面與樣式
# ==========================================
st.set_page_config(
    page_title="[診斷模式] Aurotek 自動化發布平台",
    page_icon="🔧",
    layout="wide"
)

LOGO_URL = "https://aurotek.com/wp-content/uploads/2025/07/logo.svg"
WORK_DIR = "temp_workspace"
HISTORY_FILE = "job_history.json"

# ==========================================
#              CSS (保持 v1.3 樣式)
# ==========================================
st.markdown("""
<style>
header[data-testid="stHeader"] { display: none; }
.stApp > header { display: none; }
.block-container { padding-top: 1rem !important; padding-bottom: 6rem !important; }

/* 按鈕樣式 */
[data-testid="stFileUploaderDropzoneInstructions"] > div:first-child,
[data-testid="stFileUploaderDropzoneInstructions"] > div:nth-child(2) { display: none !important; }
[data-testid="stFileUploaderDropzoneInstructions"]::before { content: "請將檔案拖放至此"; display: block; font-weight: 700; color: #31333F; }
section[data-testid="stFileUploaderDropzone"] button {
    border: 1px solid #d0d7de; background-color: #ffffff; color: transparent !important;
    position: relative; padding: 0.25rem 0.75rem; border-radius: 4px; min-height: 38px; width: auto; margin-top: 10px;
}
section[data-testid="stFileUploaderDropzone"] button::after {
    content: "瀏覽檔案"; position: absolute; color: #31333F; left: 50%; top: 50%; transform: translate(-50%, -50%); white-space: nowrap; font-weight: 500; font-size: 14px;
}
[data-testid="stFileUploaderDeleteBtn"] { border: none !important; background: transparent !important; margin-top: 0 !important; color: inherit !important; }
[data-testid="stFileUploaderDeleteBtn"]::after { content: none !important; }

/* 提示詞樣式 */
div[data-testid="stAlert"][data-style="success"], div[data-testid="stAlert"][data-style="info"] { background-color: #F0F2F6 !important; color: #31333F !important; border: 1px solid #d0d7de !important; }
div[data-testid="stAlert"] svg { color: #004280 !important; }

/* 垃圾桶與按鈕 */
div[data-testid="column"] button { border: 1px solid #eee !important; background: white !important; color: #555 !important; font-size: 0.85rem !important; min-width: 40px !important; padding: 4px 8px !important; }
div[data-testid="column"] button:hover { color: #cc0000 !important; border-color: #cc0000 !important; background: #fff5f5 !important; }
</style>
""", unsafe_allow_html=True)

# ==========================================
#              Helper Functions
# ==========================================
def cleanup_workspace():
    if os.path.exists(WORK_DIR):
        try:
            shutil.rmtree(WORK_DIR)
        except Exception as e:
            print(f"Cleanup warning: {e}")
    os.makedirs(WORK_DIR, exist_ok=True)

def reset_callback():
    cleanup_workspace()
    # 清除歷史紀錄邏輯...
    if st.session_state.get('current_file_name') and os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            if st.session_state.current_file_name in data:
                del data[st.session_state.current_file_name]
                with open(HISTORY_FILE, "w", encoding="utf-8") as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)
        except:
            pass
    st.session_state.split_jobs = []
    st.session_state.current_file_name = None
    st.session_state.ppt_meta = {"total_slides": 0, "preview_data": []}
    st.session_state.execution_results = None 
    st.session_state.reset_key += 1

def load_history(filename):
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f).get(filename, [])
        except:
            return []
    return []

def save_history(filename, jobs):
    try:
        data = {}
        if os.path.exists(HISTORY_FILE):
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                try: data = json.load(f)
                except: pass
        data[filename] = jobs
        with open(HISTORY_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except: pass

def add_split_job(total_pages):
    st.session_state.split_jobs.insert(0, {
        "id": str(uuid.uuid4())[:8], "filename": "", "start": 1, "end": total_pages,
        "category": "清潔", "subcategory": "", "client": "", "keywords": ""
    })

def remove_split_job(index):
    st.session_state.split_jobs.pop(index)

def validate_jobs(jobs, total_slides):
    errors = []
    for i, job in enumerate(jobs):
        if not job['filename'].strip(): errors.append(f"❌ 任務 {len(jobs)-i}: 檔名為空")
        if job['start'] > job['end']: errors.append(f"❌ 任務 {len(jobs)-i}: 起始頁大於結束頁")
        if job['end'] > total_slides: errors.append(f"❌ 任務 {len(jobs)-i}: 結束頁超出範圍")
    return errors

def download_file_from_url(url, dest_path):
    try:
        response = requests.get(url, stream=True, timeout=60)
        response.raise_for_status()
        with open(dest_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        return True, None
    except Exception as e:
        return False, str(e)

def copy_btn_html(text):
    return f"""<html><body><button onclick="navigator.clipboard.writeText('{text}')" style="border:1px solid #ddd;background:#fff;padding:4px 8px;border-radius:4px;cursor:pointer;">📋 複製</button></body></html>"""

# ==========================================
#              核心執行邏輯 (加強診斷)
# ==========================================
def execute_automation_logic(bot, source_path, file_prefix, jobs, auto_clean):
    main_progress = st.progress(0, text="準備開始...")
    status_area = st.empty()
    detail_bar = st.empty()
    sorted_jobs = sorted(jobs, key=lambda x: x['start'])

    # 定義回調函數
    def update_bar(text, pct):
        detail_bar.progress(pct, text=text)

    try:
        # --- Step 1 ---
        status_area.info("1️⃣ 步驟 1/5：提取 PPT 內影片並上傳至雲端...")
        main_progress.progress(5, text="Step 1: 影片雲端化")
        
        # [診斷] 檢查來源檔案
        if not os.path.exists(source_path):
            raise FileNotFoundError(f"找不到來源檔案: {source_path}")
            
        video_map = bot.extract_and_upload_videos(
            source_path,
            os.path.join(WORK_DIR, "media"),
            file_prefix=file_prefix,
            progress_callback=lambda f, c, t: update_bar(f"上傳中: {f}", c/t if t else 0),
            log_callback=lambda msg: print(f"[Step1] {msg}")
        )
        
        # [診斷] 顯示 Video Map 結果，確認 Step 1 是否真的成功
        with st.expander("🔍 [診斷] Step 1 完成，查看影片對應表 (Video Map)", expanded=True):
            st.json(video_map)
            if not video_map:
                st.warning("⚠️ 注意：沒有偵測到任何影片，若 PPT 內有影片請檢查格式。")

        # --- Step 2 ---
        status_area.info("2️⃣ 步驟 2/5：將 PPT 內的影片替換為雲端連結圖片...")
        main_progress.progress(25, text="Step 2: 連結置換")
        
        mod_path = os.path.join(WORK_DIR, "modified.pptx")
        
        # [診斷] 這是您提到最容易出錯的地方，我們包一層 try-except
        try:
            bot.replace_videos_with_images(
                source_path,
                mod_path,
                video_map,
                progress_callback=lambda c, t: update_bar(f"處理投影片 {c}/{t}", c/t if t else 0)
            )
        except Exception as e_step2:
            st.sidebar.error(f"❌ 錯誤發生在步驟二 (影片替換)！\n\n原因: {e_step2}")
            st.sidebar.code(traceback.format_exc())
            raise e_step2  # 拋出錯誤讓主流程中止

        # --- Step 3 ---
        status_area.info("3️⃣ 步驟 3/5：進行檔案壓縮與瘦身...")
        main_progress.progress(45, text="Step 3: 檔案瘦身")
        slim_path = os.path.join(WORK_DIR, "slim.pptx")
        bot.shrink_pptx(mod_path, slim_path, progress_callback=lambda c, t: update_bar("壓縮中...", c/t if t else 0))

        # --- Step 4 ---
        status_area.info("4️⃣ 步驟 4/5：拆分並上傳至 Google Slides...")
        main_progress.progress(65, text="Step 4: 拆分發布")
        results = bot.split_and_upload(
            slim_path, sorted_jobs, file_prefix,
            progress_callback=lambda f, c, t: update_bar(f"上傳簡報: {f}", c/t if t else 0),
            log_callback=print
        )

        # 檢查是否有過大檔案
        oversized = [r for r in results if r.get('error_too_large')]
        if oversized:
            st.error("⛔️ 檔案過大，無法轉換。")
            return

        # --- Step 5 ---
        status_area.info("5️⃣ 步驟 5/5：優化線上播放器...")
        main_progress.progress(85, text="Step 5: 內嵌優化")
        final_results = bot.embed_videos_in_slides(results, progress_callback=lambda c, t: update_bar("優化中...", c/t if t else 0), log_callback=print)

        # --- Final ---
        status_area.info("📝 最後步驟：寫入資料庫...")
        main_progress.progress(95, text="Final: 寫入資料庫")
        bot.log_to_sheets(final_results, log_callback=print)

        main_progress.progress(100, text="任務完成")
        status_area.info("**成功：** 所有自動化流程執行完畢。", icon=None)
        
        if auto_clean: cleanup_workspace()
        
        st.session_state.execution_results = {"results": final_results, "prefix": file_prefix}

    except Exception as e:
        # [關鍵] 捕捉所有錯誤並顯示在側邊欄
        st.sidebar.error("❌ 執行流程發生錯誤！請截圖此畫面給開發者。")
        st.sidebar.error(f"錯誤類型: {type(e).__name__}")
        st.sidebar.error(f"錯誤訊息: {str(e)}")
        with st.sidebar.expander("查看完整程式碼追蹤 (Traceback)", expanded=True):
            st.code(traceback.format_exc())
        # 同時在主畫面顯示
        st.error("程式發生錯誤，請查看左側邊欄的詳細資訊。")

# ==========================================
#              主程式介面邏輯
# ==========================================
os.makedirs(WORK_DIR, exist_ok=True)

# 狀態初始化
if 'split_jobs' not in st.session_state: st.session_state.split_jobs = []
if 'reset_key' not in st.session_state: st.session_state.reset_key = 0
if 'execution_results' not in st.session_state: st.session_state.execution_results = None
if 'bot' not in st.session_state:
    try:
        bot_instance = PPTAutomationBot()
        if bot_instance.creds: st.session_state.bot = bot_instance
    except: pass

if 'current_file_name' not in st.session_state: st.session_state.current_file_name = None
if 'ppt_meta' not in st.session_state: st.session_state.ppt_meta = {"total_slides": 0, "preview_data": []}

# UI Header
components.html(f"""<div style="width:100%;display:flex;flex-direction:column;align-items:center;margin:4px 0 2px 0;"><img src="{LOGO_URL}" style="width:300px;"><div style="margin-top:4px;color:gray;font-size:1rem;letter-spacing:2px;">簡報案例自動化發布平台</div></div>""", height=78)
st.info("功能說明： 上傳PPT → 線上拆分 → 影片雲端化 → 內嵌優化 → 簡報雲端化 → 寫入和椿資料庫")

# 機器人檢查
if 'bot' not in st.session_state:
    st.error("❌ 機器人未初始化 (憑證錯誤)，請檢查 Secrets。")

# Step 1
with st.container(border=True):
    st.subheader("步驟一：選擇檔案來源")
    input_method = st.radio("上傳方式", ["本地檔案", "線上檔案"], horizontal=True)
    uploaded_file = None
    source_path = os.path.join(WORK_DIR, "source.pptx")
    file_name_for_logic = None
    
    if input_method == "本地檔案":
        uploaded_file = st.file_uploader("請選擇 PPTX 檔案", type=['pptx'], label_visibility="collapsed", key=f"uploader_{st.session_state.reset_key}")
        if uploaded_file:
            file_name_for_logic = uploaded_file.name
            if st.session_state.current_file_name != file_name_for_logic:
                cleanup_workspace()
                with open(source_path, "wb") as f: f.write(uploaded_file.getbuffer())
            elif not os.path.exists(source_path):
                 with open(source_path, "wb") as f: f.write(uploaded_file.getbuffer())
    else:
        url_input = st.text_input("請輸入 PPTX 下載網址", key=f"url_{st.session_state.reset_key}")
        if url_input and st.button("下載"):
            cleanup_workspace()
            success, err = download_file_from_url(url_input, source_path)
            if success:
                file_name_for_logic = "downloaded.pptx"
                st.info("下載成功", icon="✅")
            else:
                st.error(f"下載失敗: {err}")

    if file_name_for_logic and os.path.exists(source_path):
        if st.session_state.current_file_name != file_name_for_logic:
            try:
                prs = Presentation(source_path)
                st.session_state.ppt_meta["total_slides"] = len(prs.slides)
                st.session_state.ppt_meta["preview_data"] = [{"頁碼": i+1} for i in range(len(prs.slides))]
                st.session_state.current_file_name = file_name_for_logic
                st.session_state.split_jobs = load_history(file_name_for_logic) or []
                st.session_state.execution_results = None
                st.info(f"**已讀取：** {file_name_for_logic} (共 {len(prs.slides)} 頁)", icon=None)
            except Exception as e:
                st.error(f"檔案讀取失敗: {e}")
                st.session_state.current_file_name = None

# Step 2 & 3
if st.session_state.current_file_name:
    with st.expander("👁️ 查看頁碼對照表"):
        st.dataframe(st.session_state.ppt_meta["preview_data"], use_container_width=True)

    with st.container(border=True):
        c1, c2 = st.columns([3, 1])
        c1.subheader("步驟二：設定拆分任務")
        if c2.button("➕ 新增任務", type="primary", use_container_width=True):
            add_split_job(st.session_state.ppt_meta["total_slides"])

        if not st.session_state.split_jobs:
            st.info("尚未建立任務，請點擊上方按鈕新增。")
        
        for i, job in enumerate(st.session_state.split_jobs):
            with st.container(border=True):
                c_title, c_del = st.columns([0.85, 0.15])
                c_title.markdown(f"**任務 {len(st.session_state.split_jobs)-i}**")
                if c_del.button("🗑️ 刪除", key=f"del_{job['id']}"):
                    remove_split_job(i)
                    st.rerun()
                
                c_a, c_b, c_c = st.columns([3, 1.5, 1.5])
                job["filename"] = c_a.text_input("檔名", value=job["filename"], key=f"f_{job['id']}")
                job["start"] = c_b.number_input("起始", 1, st.session_state.ppt_meta["total_slides"], job["start"], key=f"s_{job['id']}")
                job["end"] = c_c.number_input("結束", 1, st.session_state.ppt_meta["total_slides"], job["end"], key=f"e_{job['id']}")
                
                c_d, c_e, c_f, c_g = st.columns(4)
                job["category"] = c_d.selectbox("類型", ["清潔", "配送", "購物", "AURO"], key=f"cat_{job['id']}")
                job["subcategory"] = c_e.text_input("子分類", value=job["subcategory"], key=f"sub_{job['id']}")
                job["client"] = c_f.text_input("客戶", value=job["client"], key=f"cli_{job['id']}")
                job["keywords"] = c_g.text_input("關鍵字", value=job["keywords"], key=f"key_{job['id']}")
        
        save_history(st.session_state.current_file_name, st.session_state.split_jobs)

    # Step 3 (Only if jobs exist)
    if st.session_state.split_jobs:
        with st.container(border=True):
            st.subheader("步驟三：執行任務")
            auto_clean = st.checkbox("任務完成後自動清除暫存檔", value=True)
            if st.button("執行雲端化任務", type="primary", use_container_width=True):
                errs = validate_jobs(st.session_state.split_jobs, st.session_state.ppt_meta["total_slides"])
                if errs:
                    for e in errs: st.error(e)
                else:
                    if st.session_state.bot:
                        execute_automation_logic(
                            st.session_state.bot,
                            os.path.join(WORK_DIR, "source.pptx"),
                            os.path.splitext(st.session_state.current_file_name)[0],
                            st.session_state.split_jobs,
                            auto_clean
                        )
                        st.rerun()
                    else:
                        st.error("Bot 未初始化")

# Step 4 & Footer
if st.session_state.execution_results:
    st.markdown("<div id='step4-anchor'></div>", unsafe_allow_html=True)
    with st.container(border=True):
        st.subheader("步驟四：產出結果")
        results = st.session_state.execution_results["results"]
        pfx = st.session_state.execution_results["prefix"]
        
        # 簡單表格渲染
        rows = ""
        for r in results:
            if 'final_link' in r:
                rows += f"""<tr style="border-bottom:1px solid #eee;"><td style="padding:8px;">[{pfx}]_{r['filename']}</td><td style="padding:8px;"><a href="{r['final_link']}" target="_blank">開啟</a></td></tr>"""
        
        if rows:
            st.markdown(f"""<table style="width:100%;font-size:14px;"><tr><th style="text-align:left;padding:8px;">檔案</th><th style="padding:8px;">連結</th></tr>{rows}</table>""", unsafe_allow_html=True)
        else:
            st.warning("無結果")
            
    components.html("""<script>setTimeout(function(){try{window.parent.document.getElementById('step4-anchor').scrollIntoView({behavior:'smooth',block:'start'});}catch(e){}},500);</script>""", height=0)

if st.session_state.current_file_name:
    st.markdown("<div style='margin-top: 40px;'></div>", unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    with c1:
        st.button("清除任務，上傳新簡報", type="primary", on_click=reset_callback, use_container_width=True)
    with c2:
<<<<<<< HEAD
        st.link_button("前往「和椿數位資源庫」", "https://aurotek.pse.is/puducases", type="primary", use_container_width=True)
=======
        st.link_button("前往「和椿數位資源庫」", "https://aurotek.pse.is/puducases", type="primary", use_container_width=True)
>>>>>>> 6b3f31d (Update app.py to v1.3 with UI fixes)
